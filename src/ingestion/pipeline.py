"""
Document ingestion pipeline with incremental indexing.

Supported formats
-----------------
  .pdf          docling (rich: tables, figures, structure) → pypdf fallback
  .docx / .odt  docling → python-docx fallback
  .pptx         docling
  .xlsx / .xls  docling → openpyxl fallback (one row per Document)
  .html / .htm  docling → regex tag-strip fallback
  .epub         docling
  .xml          lxml → regex fallback
  .rtf          striprtf → plaintext fallback
  .json         flattened key-value text
  .txt / .md    direct UTF-8 read
  .csv          one Document per row (pandas → csv.DictReader fallback)
  .png .jpg .jpeg .tiff .bmp .webp
                docling OCR (tesseract required) → skipped if unavailable

Incremental strategy
--------------------
  Each file is hashed (SHA-256 of raw bytes).  Already-seen hashes are
  skipped entirely — no re-embedding, no re-indexing on re-run.
  The hash store persists to disk so cold restarts are also incremental.
"""
import csv
import hashlib
import json
import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterator, List

logger = logging.getLogger(__name__)

SUPPORTED = {
    # document
    ".pdf", ".docx", ".odt", ".pptx", ".epub",
    # spreadsheet
    ".xlsx", ".xls",
    # web / markup
    ".html", ".htm", ".xml",
    # plain text
    ".txt", ".md", ".rtf",
    # structured data
    ".csv", ".json",
    # images (requires docling + tesseract)
    ".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp",
}


# ── data model ────────────────────────────────────────────────────────────────

@dataclass
class Document:
    doc_id: str
    source: str
    text: str
    metadata: dict = field(default_factory=dict)


# ── text normalisation ────────────────────────────────────────────────────────

def _norm(t: str) -> str:
    if not t:
        return ""
    t = t.replace("\x00", "")
    t = re.sub(r"\n{3,}", "\n\n", t)
    return re.sub(r" {2,}", " ", t).strip()


# ── docling universal loader ──────────────────────────────────────────────────

def _load_via_docling(path: Path, with_ocr: bool = False) -> str:
    """
    Primary loader for most formats.
    Returns Markdown-formatted text preserving structure (headings, tables, lists).
    Returns None on ImportError/failure so callers can apply a format-specific fallback.
    """
    try:
        from docling.document_converter import DocumentConverter
        from docling.datamodel.pipeline_options import PdfPipelineOptions
        from docling.datamodel.base_models import InputFormat
        from docling.document_converter import PdfFormatOption

        opts = PdfPipelineOptions()
        opts.do_ocr = with_ocr
        opts.do_table_structure = True

        converter = DocumentConverter(
            format_options={InputFormat.PDF: PdfFormatOption(pipeline_options=opts)}
        )
        result = converter.convert(str(path))
        text = result.document.export_to_markdown()
        logger.debug("docling OK: %s (%d chars)", path.name, len(text))
        return text or None
    except ImportError:
        logger.debug("docling not installed — fallback for %s", path.name)
        return None
    except Exception as e:
        logger.debug("docling failed for %s: %s — fallback", path.name, e)
        return None


# ── per-format fallback loaders ───────────────────────────────────────────────

def _load_pdf(path: Path) -> str:
    t = _load_via_docling(path)
    if t:
        return t
    try:
        from pypdf import PdfReader
        pages = PdfReader(str(path)).pages
        return "\n".join(p.extract_text() or "" for p in pages)
    except Exception as e:
        logger.warning("PDF load failed %s: %s", path.name, e)
        return ""


def _load_docx(path: Path) -> str:
    t = _load_via_docling(path)
    if t:
        return t
    try:
        from docx import Document as DocxDoc
        doc = DocxDoc(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        logger.warning("DOCX load failed %s: %s", path.name, e)
        return ""


def _load_pptx(path: Path) -> str:
    t = _load_via_docling(path)
    if t:
        return t
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts)
    except Exception as e:
        logger.warning("PPTX load failed %s: %s", path.name, e)
        return ""


def _load_epub(path: Path) -> str:
    t = _load_via_docling(path)
    if t:
        return t
    try:
        import zipfile, re as _re
        with zipfile.ZipFile(str(path)) as z:
            texts = []
            for name in z.namelist():
                if name.endswith((".html", ".xhtml")):
                    raw = z.read(name).decode("utf-8", errors="replace")
                    texts.append(_re.sub(r"<[^>]+>", " ", raw))
        return " ".join(texts)
    except Exception as e:
        logger.warning("EPUB load failed %s: %s", path.name, e)
        return ""


def _load_html(path: Path) -> str:
    t = _load_via_docling(path)
    if t:
        return t
    raw = path.read_text(encoding="utf-8", errors="replace")
    return re.sub(r"<[^>]+>", " ", raw)


def _load_xml(path: Path) -> str:
    try:
        import lxml.etree as ET
        tree = ET.parse(str(path))
        return " ".join(tree.getroot().itertext())
    except ImportError:
        pass
    except Exception as e:
        logger.warning("lxml XML parse failed %s: %s", path.name, e)
    # regex fallback
    raw = path.read_text(encoding="utf-8", errors="replace")
    return re.sub(r"<[^>]+>", " ", raw)


def _load_rtf(path: Path) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        return rtf_to_text(path.read_text(encoding="utf-8", errors="replace"))
    except ImportError:
        pass
    # crude fallback: strip RTF control words
    raw = path.read_text(encoding="utf-8", errors="replace")
    return re.sub(r"\\[a-zA-Z0-9]+|[{}]", " ", raw)


def _load_image(path: Path) -> str:
    """OCR via docling (requires tesseract). Returns empty if unavailable."""
    t = _load_via_docling(path, with_ocr=True)
    if t:
        return t
    logger.info("Image OCR unavailable — skipping %s (install docling + tesseract)", path.name)
    return ""


def _load_json(path: Path) -> str:
    """Flatten JSON to key: value lines for retrieval."""
    try:
        data = json.loads(path.read_text(encoding="utf-8", errors="replace"))
    except Exception as e:
        logger.warning("JSON parse failed %s: %s", path.name, e)
        return path.read_text(encoding="utf-8", errors="replace")

    def _flatten(obj, prefix=""):
        parts = []
        if isinstance(obj, dict):
            for k, v in obj.items():
                parts.extend(_flatten(v, f"{prefix}{k}: " if not prefix else f"{prefix}.{k}: "))
        elif isinstance(obj, list):
            for i, v in enumerate(obj):
                parts.extend(_flatten(v, f"{prefix}[{i}]: "))
        else:
            if str(obj).strip():
                parts.append(f"{prefix}{obj}")
        return parts

    return "\n".join(_flatten(data))


def _load_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


# ── Excel: multi-sheet, row-per-Document ─────────────────────────────────────

def _excel_rows(path: Path) -> Iterator[str]:
    """Yield one text string per row across all sheets."""
    # try docling first (returns full markdown with tables)
    t = _load_via_docling(path)
    if t:
        yield t
        return
    # openpyxl fallback
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        for sheet in wb.worksheets:
            headers = None
            for i, row in enumerate(sheet.iter_rows(values_only=True)):
                if not any(c is not None for c in row):
                    continue
                if i == 0:
                    headers = [str(c).strip() if c is not None else f"col{j}"
                               for j, c in enumerate(row)]
                else:
                    if headers:
                        pairs = [f"{h}: {v}" for h, v in zip(headers, row)
                                 if v is not None and str(v).strip()]
                        if pairs:
                            yield " | ".join(pairs)
                    else:
                        row_vals = [str(v).strip() for v in row
                                    if v is not None and str(v).strip()]
                        if row_vals:
                            yield " | ".join(row_vals)
        return
    except ImportError:
        pass
    except Exception as e:
        logger.warning("openpyxl failed %s: %s", path.name, e)
    # pandas fallback
    try:
        import pandas as pd
        all_sheets = pd.read_excel(str(path), sheet_name=None, dtype=str)
        for _sheet_name, df in all_sheets.items():
            df = df.fillna("")
            for _, row in df.iterrows():
                pairs = [f"{col}: {val}" for col, val in row.items() if str(val).strip()]
                if pairs:
                    yield " | ".join(pairs)
    except Exception as e:
        logger.warning("pandas Excel fallback failed %s: %s", path.name, e)


def _csv_rows(path: Path) -> Iterator[str]:
    try:
        import pandas as pd
        df = pd.read_csv(str(path), dtype=str).fillna("")
        for _, row in df.iterrows():
            pairs = [f"{col}: {val}" for col, val in row.items() if str(val).strip()]
            if pairs:
                yield " | ".join(pairs)
        return
    except ImportError:
        pass
    except Exception as e:
        logger.debug("pandas CSV failed, using stdlib: %s", e)
    # stdlib fallback
    try:
        with path.open(newline="", encoding="utf-8", errors="replace") as f:
            for row in csv.DictReader(f):
                text = " | ".join(f"{k}: {v}" for k, v in row.items() if v)
                if text:
                    yield text
    except Exception as e:
        logger.warning("CSV failed %s: %s", path.name, e)


# ── loader dispatch table ─────────────────────────────────────────────────────

_SINGLE_LOADERS = {
    ".pdf":  _load_pdf,
    ".docx": _load_docx,
    ".odt":  _load_docx,    # docling handles ODT
    ".pptx": _load_pptx,
    ".epub": _load_epub,
    ".html": _load_html,
    ".htm":  _load_html,
    ".xml":  _load_xml,
    ".rtf":  _load_rtf,
    ".txt":  _load_text,
    ".md":   _load_text,
    ".json": _load_json,
    ".png":  _load_image,
    ".jpg":  _load_image,
    ".jpeg": _load_image,
    ".tiff": _load_image,
    ".tif":  _load_image,
    ".bmp":  _load_image,
    ".webp": _load_image,
}

# Formats that produce one Document per row/entry
_ROW_LOADERS = {
    ".csv":  _csv_rows,
    ".xlsx": _excel_rows,
    ".xls":  _excel_rows,
}


# ── hash store ────────────────────────────────────────────────────────────────

class HashStore:
    """Persist SHA-256 hashes to disk for incremental ingestion."""

    def __init__(self, path):
        self.path = Path(path)
        self._s: dict = json.loads(self.path.read_text()) if self.path.exists() else {}

    def seen(self, did: str) -> bool:
        return did in self._s

    def add(self, did: str, src: str):
        self._s[did] = src

    def remove(self, did: str):
        self._s.pop(did, None)

    def save(self):
        self.path.parent.mkdir(parents=True, exist_ok=True)
        self.path.write_text(json.dumps(self._s, indent=2))

    def all_ids(self) -> set:
        return set(self._s)

    def count(self) -> int:
        return len(self._s)


# ── pipeline ──────────────────────────────────────────────────────────────────

class IngestPipeline:
    def __init__(self, hash_store: HashStore):
        self.hs = hash_store

    def _hash(self, p: Path) -> str:
        return hashlib.sha256(p.read_bytes()).hexdigest()

    def ingest_file(self, path) -> List[Document]:
        path = Path(path)
        suffix = path.suffix.lower()
        if suffix not in SUPPORTED:
            logger.debug("Unsupported format — skipped: %s", path.name)
            return []

        did = self._hash(path)
        if self.hs.seen(did):
            logger.debug("Unchanged (skipped): %s", path.name)
            return []

        meta = {"filename": path.name, "suffix": suffix,
                "size_bytes": path.stat().st_size}

        # ── row-based formats (CSV, Excel) ────────────────────────────────
        if suffix in _ROW_LOADERS:
            row_iter = _ROW_LOADERS[suffix](path)
            docs = [
                Document(f"{did}_r{i}", str(path), _norm(row), {**meta, "row": i})
                for i, row in enumerate(row_iter)
                if _norm(row)
            ]
            if docs:
                self.hs.add(did, str(path))
            logger.info("Ingested %s (%s) → %d row-docs",
                        path.name, suffix, len(docs))
            return docs

        # ── single-document formats ───────────────────────────────────────
        loader = _SINGLE_LOADERS.get(suffix)
        if loader is None:
            return []
        raw = loader(path)
        text = _norm(raw)
        if not text:
            logger.warning("Empty text from %s — skipped", path.name)
            return []
        self.hs.add(did, str(path))
        logger.info("Ingested %s (%s, %d chars)", path.name, suffix, len(text))
        return [Document(did, str(path), text, meta)]

    def ingest_dir(self, directory, recursive: bool = True) -> List[Document]:
        directory = Path(directory)
        pattern = "**/*" if recursive else "*"
        all_docs, skipped = [], 0
        for f in sorted(directory.glob(pattern)):
            if not f.is_file():
                continue
            docs = self.ingest_file(f)
            if not docs:
                skipped += 1
            all_docs.extend(docs)
        self.hs.save()
        logger.info(
            "Ingest complete — %d new doc(s), %d unchanged/unsupported skipped",
            len(all_docs), skipped,
        )
        return all_docs

    @staticmethod
    def supported_extensions() -> List[str]:
        return sorted(SUPPORTED)
